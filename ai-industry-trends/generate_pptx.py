# -*- coding: utf-8 -*-
"""Generate AI Industry Trends PPTX presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
NAVY     = RGBColor(0x0A, 0x1E, 0x46)
BLUE     = RGBColor(0x1E, 0x5A, 0xC4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
GRAY     = RGBColor(0x80, 0x90, 0xA8)
DARK     = RGBColor(0x3A, 0x4A, 0x6A)
GRID_C   = RGBColor(0xE8, 0xEC, 0xF2)
CARD_BDR = RGBColor(0xCC, 0xD2, 0xDD)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def tb(slide, left, top, width, height, text, font_size=18, color=NAVY,
       bold=False, font_name='Arial', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Simple single-paragraph text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def mono(slide, left, top, width, height, text, font_size=12, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    """Monospace text box."""
    return tb(slide, left, top, width, height, text, font_size, color, bold, 'Consolas', align)

def center_tb(slide, top, height, text, font_size=18, color=NAVY, bold=False, font_name='Arial', align=PP_ALIGN.CENTER):
    """Full-width centered text box."""
    return tb(slide, Inches(1), top, W - Inches(2), height, text, font_size, color, bold, font_name, align)

def add_section_slide(slide, num_str, title, subtitle):
    add_bg(slide, WHITE)
    # Grid
    for i in range(0, 40):
        x = Emu(i * 40 * 12700)
        add_rect(slide, x, Inches(0), Pt(0.5), H, fill_color=GRID_C)
        y = Emu(i * 40 * 12700)
        add_rect(slide, Inches(0), y, W, Pt(0.5), fill_color=GRID_C)
    center_tb(slide, Inches(0.5), Inches(3.5), num_str, 160, GRID_C, True, 'Consolas', PP_ALIGN.CENTER)
    mono(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(0.4),
         'CHAPTER ' + num_str, 13, BLUE, True)
    tb(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(1.0),
       title, 48, NAVY, True)
    tb(slide, Inches(1.2), Inches(3.8), Inches(11), Inches(0.6),
       subtitle, 17, DARK)
    add_rect(slide, Inches(1.2), Inches(3.6), Inches(1.5), Pt(4), fill_color=BLUE)

def add_card(slide, left, top, w, h, title, body, accent=None):
    add_rect(slide, left, top, w, h, fill_color=WHITE, border_color=NAVY, border_width=Pt(1.5))
    tb(slide, left + Inches(0.3), top + Inches(0.2), w - Inches(0.6), Inches(0.45),
       title, 17, NAVY, True)
    tb(slide, left + Inches(0.3), top + Inches(0.7), w - Inches(0.6), h - Inches(1.3),
       body, 12, DARK)
    if accent:
        mono(slide, left + Inches(0.3), top + h - Inches(0.45), w - Inches(0.6), Inches(0.35),
             accent, 10, BLUE, True)

def add_insight(slide, left, top, w, h, label, text):
    add_rect(slide, left, top, Pt(5), h, fill_color=BLUE)
    mono(slide, left + Inches(0.25), top + Inches(0.05), w - Inches(0.25), Inches(0.3),
         label, 11, BLUE, True)
    tb(slide, left + Inches(0.25), top + Inches(0.38), w - Inches(0.25), h - Inches(0.45),
       text, 13, NAVY)

def add_sn(slide, num, total=16):
    mono(slide, W - Inches(1.8), H - Inches(0.6), Inches(1.5), Inches(0.35),
         f'{num} / {total}', 9, GRAY, align=PP_ALIGN.RIGHT)

def add_footer(slide):
    mono(slide, Inches(0.8), H - Inches(0.6), Inches(8), Inches(0.35),
         'AI Industry Trends 2026', 9, GRAY)
    add_rect(slide, Inches(0.8), H - Inches(0.5), W - Inches(1.6), Pt(1), fill_color=NAVY)

def add_stat_row(slide, top, stats):
    """stats: list of (value, label, is_accent)"""
    for i, (val, lbl, is_accent) in enumerate(stats):
        left = Inches(1 + i * 3.0)
        add_rect(slide, left, top, Inches(2.6), Inches(1.15),
                 fill_color=WHITE, border_color=CARD_BDR, border_width=Pt(1))
        c = BLUE if is_accent else NAVY
        tb(slide, left + Inches(0.15), top + Inches(0.1), Inches(2.3), Inches(0.55),
           val, 30, c, True, 'Consolas', PP_ALIGN.CENTER)
        tb(slide, left + Inches(0.15), top + Inches(0.6), Inches(2.3), Inches(0.4),
           lbl, 11, GRAY, align=PP_ALIGN.CENTER)

def add_table(slide, left, top, w, data, col_widths):
    """data: list of lists; first row = header"""
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.42 * rows)).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.text = data[r][c]
            p.font.size = Pt(12)
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = BLUE
                p.font.size = Pt(11)
                p.font.name = 'Consolas'
            elif c == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY
            else:
                p.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if r == 0 else WHITE

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
# Grid pattern
for i in range(0, 40):
    x = Emu(i * 40 * 12700)
    add_rect(sl, x, Inches(0), Pt(0.5), H, fill_color=GRID_C)
    y = Emu(i * 40 * 12700)
    add_rect(sl, Inches(0), y, W, Pt(0.5), fill_color=GRID_C)

mono(sl, Inches(1.2), Inches(1.4), Inches(10), Inches(0.5),
     '◆ AI INFRASTRUCTURE · 2026', 13, BLUE, True)
tb(sl, Inches(1.2), Inches(2.1), Inches(10.5), Inches(1.5),
   'AI 产业发展趋势', 64, NAVY, True)
tb(sl, Inches(1.2), Inches(3.2), Inches(10.5), Inches(0.8),
   '从算力到互联的硬科技革命', 36, BLUE, False)
tb(sl, Inches(1.2), Inches(4.3), Inches(9), Inches(0.9),
   'CPU · GPU · 光模块 · PCB · 存储 · 互联 · 散热 ——'
   ' 拆解 AI 基础设施的每一个关键环节，'
   '看清技术演进与投资脉络。',
   16, DARK)
mono(sl, Inches(1.2), Inches(5.4), Inches(10), Inches(0.5),
     '2026 年 6 月    |    AI 产业研究团队    |    ~30 min', 12, GRAY)
add_sn(sl, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1.2), Inches(0.6), Inches(10), Inches(0.5), '# AGENDA', 13, BLUE, True)
tb(sl, Inches(1.2), Inches(1.1), Inches(10), Inches(0.9),
   '今天的内容路线图', 42, NAVY, True)

agenda = [
    ('01', 'CPU：算力基石', '~3 min'),
    ('02', 'GPU：训练与推理引擎', '~3 min'),
    ('03', '光模块：数据洪流通道', '~3 min'),
    ('04', 'PCB：高速信号载体', '~2 min'),
    ('05', '存储：大模型记忆基石', '~3 min'),
    ('06', '互联与网络', '~2 min'),
    ('07', '散热与能效', '~2 min'),
    ('08', '市场格局与展望', '~4 min'),
]
for i, (num, title, dur) in enumerate(agenda):
    col, row = i % 2, i // 2
    left = Inches(1.2 + col * 5.6)
    top  = Inches(2.3 + row * 1.05)
    add_rect(sl, left, top, Inches(5.2), Inches(0.9),
             fill_color=WHITE, border_color=CARD_BDR, border_width=Pt(1))
    mono(sl, left + Inches(0.25), top + Inches(0.15), Inches(0.6), Inches(0.55),
         num, 24, BLUE, True)
    tb(sl, left + Inches(0.9), top + Inches(0.2), Inches(3.2), Inches(0.5),
       title, 16, NAVY, True)
    mono(sl, left + Inches(4.2), top + Inches(0.22), Inches(1.0), Inches(0.45),
         dur, 11, GRAY, align=PP_ALIGN.RIGHT)
add_sn(sl, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — CPU SECTION
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_slide(sl, '01', 'CPU：算力基石',
                  '从通用计算到 AI 专用 —— 服务器 CPU 的新战场')
add_sn(sl, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — CPU DETAIL
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# CPU · 算力基石', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   'x86 与 ARM 两强争霸，AI 专用加速单元成为标配', 36, NAVY, True)

cpu_cards = [
    ('\U0001f535 Intel Xeon 6',
     'Granite Rapids · 最多 128 核 · 内置 AMX 矩阵引擎，'
     'BF16/INT8 推理加速。单路 12 通道 DDR5，PCIe 5.0 × 88 lane。',
     'AMX · 2048-bit'),
    ('\U0001f534 AMD EPYC 9005',
     'Turin · Zen 5 架构 · 最多 192 核 384 线程 · AVX-512 全吞吐。'
     '12 通道 DDR5-6000，CXL 2.0 内存扩展。',
     '192 Cores · 500W TDP'),
    ('\U0001f7e2 ARM Neoverse V3 / N3',
     'AWS Graviton4 · NVIDIA Grace · 每瓦性能领先 x86 约 30-50%。'
     'AmpereOne 192 核云原生 CPU 已量产。',
     'ARM v9.2 · SVE2'),
]
for i, (title, body, accent) in enumerate(cpu_cards):
    add_card(sl, Inches(1 + i * 3.85), Inches(2.3), Inches(3.55), Inches(2.9), title, body, accent)

add_insight(sl, Inches(1), Inches(5.5), Inches(11.5), Inches(1.2),
            '\U0001f4a1 关键趋势',
            'CPU 不再是「通用计算」的代名词。'
            '内置 AI 加速器（AMX、SVE、NPU）正让 CPU 在推理侧重新夺回话语权。'
            'ARM 阵营凭借每瓦性能优势在云原生场景加速渗透，'
            'x86 以生态和单核峰值守住基本盘。')
add_sn(sl, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — GPU SECTION
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_slide(sl, '02', 'GPU：训练与推理的算力引擎',
                  '从 Blackwell 到 ASIC —— 算力军备竞赛的核心战场')
add_sn(sl, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — GPU DETAIL
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# GPU · 算力引擎', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   'Blackwell 领跑，AMD + 自研 ASIC 双线追击', 36, NAVY, True)

add_stat_row(sl, Inches(2.2), [
    ('20K+', 'B200 FP4 TFLOPS', True),
    ('8.3K', 'MI450X FP8 TFLOPS', False),
    ('192 GB', 'B200 HBM3e 显存', False),
    ('8 TB/s', 'HBM3e 带宽', False),
])

gpu_data = [
    ['', 'NVIDIA B200', 'AMD MI450X', 'Google TPU v6', '华为 Ascend 910C'],
    ['制程', 'TSMC 4NP', 'TSMC N3', 'TSMC N3', 'SMIC N+2'],
    ['FP16 TFLOPS', '4,500 ▲', '2,300 ▲', '1,800 ▲', '800 ▲'],
    ['显存', '192 GB HBM3e', '288 GB HBM3e', '96 GB HBM3', '128 GB HBM2e'],
    ['TDP', '1,000W', '750W', '600W', '550W'],
]
add_table(sl, Inches(1), Inches(3.6), Inches(11.3), gpu_data,
          [Inches(1.6), Inches(2.425), Inches(2.425), Inches(2.425), Inches(2.425)])

add_insight(sl, Inches(1), Inches(5.9), Inches(11.5), Inches(1.0),
            '\U0001f4a1 关键趋势',
            'NVIDIA 以 Blackwell 架构 + NVLink 5 + 液冷整机柜方案建立系统级壁垒。'
            'AMD 以更大显存和开放生态（ROCm 6.0）追赶。'
            'Google / AWS / Meta 自研 ASIC 推理芯片加速部署。')
add_sn(sl, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — OPTICAL SECTION
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_slide(sl, '03', '光模块：数据洪流的「光速」通道',
                  '800G → 1.6T → 3.2T，光电互联的速度革命')
add_sn(sl, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — OPTICAL DETAIL
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# 光模块 · 互联通道', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   '800G 大规模出货，1.6T 即将商用，硅光技术拐点已至', 36, NAVY, True)

add_card(sl, Inches(1), Inches(2.2), Inches(5.5), Inches(2.3),
         '\U0001f537 800G OSFP / QSFP-DD800',
         '2026 年 AI 训练集群标配。DR8 / FR4 / SR8 多模方案齐全。'
         '单模块功耗 ~14W，硅光版本降至 ~10W。'
         '主要供应商：中际旭创、Coherent、Finisar。',
         '8×100G PAM4')
add_card(sl, Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.3),
         '\U0001f536 1.6T / 3.2T 下一代',
         '1.6T OSFP-XD 预计 2026H2 量产，200G/lane × 8。'
         'CPO（共封装光学）方案将光引擎集成至 Switch ASIC 旁，'
         '功耗降低 40%+。',
         '200G/lane · CPO')

add_card(sl, Inches(1), Inches(4.7), Inches(3.55), Inches(1.6),
         'LPO（线性驱动）',
         '去掉 DSP 芯片，功耗降低 50%，延迟 < 5ns。2026 年 800G LPO 开始小批量部署。')
add_card(sl, Inches(4.85), Inches(4.7), Inches(3.55), Inches(1.6),
         '硅光（SiPh）',
         'CMOS 工艺兼容，低成本大规模制造。Intel / Ayar Labs / 曦智科技领跑。'
         '2027 渗透率预计 > 30%。')
add_card(sl, Inches(8.7), Inches(4.7), Inches(3.55), Inches(1.6),
         '薄膜铌酸锂（TFLN）',
         '超高带宽调制器，200G+ 单波速率关键使能技术。'
         'HyperLight / 光库科技布局。')

add_insight(sl, Inches(1), Inches(6.5), Inches(11.5), Inches(0.7),
            '\U0001f4a1 关键趋势',
            '光模块速率每 2 年翻一番。'
            'LPO + 硅光两条技术路线并行演进，'
            '产业链话语权向模块厂商和硅光代工厂集中。')
add_sn(sl, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — PCB SECTION
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_slide(sl, '04', 'PCB：高速信号的高密度载体',
                  '从 M8 到 M9 材料，AI 服务器 PCB 的价值重构')
add_sn(sl, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — PCB DETAIL
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# PCB · 信号载体', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   'AI 服务器 PCB 四大升级方向', 36, NAVY, True)

pcb_cards = [
    ('\U0001f4d0 层数升级：从 20L → 30L+',
     'GPU 模组板（UBB）普遍 26-30 层，Switch 背板 34 层+。'
     '高层数带来更高的布线密度和更复杂的电源分配网络（PDN）。',
     '30L+ · UBB'),
    ('\U0001f527 材料升级：M7 → M8 → M9',
     '超低损耗 CCL 需求爆发。松下 Megtron 8、台耀 TU-983 为主流方案。'
     'M9 级 Df < 0.0015 @ 10GHz。',
     'Df < 0.0015 · M9'),
    ('\U0001f4cf HDI + Any-Layer 工艺',
     'NVIDIA GB300 NVL72 整机柜采用超大尺寸 HDI 板，'
     '单板面积 700×600mm+。Any-Layer 技术从手机走向数据中心。',
     '700×600mm · Any-Layer'),
    ('\U0001f4aa 供需格局：高端产能紧缺',
     '全球高端 PCB 产能集中在臻鼎、欣兴、Ibiden、三星电机。'
     'AI 需求拉动高端 PCB 市场 CAGR 25%+（2024-2030）。',
     'CAGR 25% · 供不应求'),
]
for i, (title, body, accent) in enumerate(pcb_cards):
    col, row = i % 2, i // 2
    add_card(sl, Inches(1 + col * 5.85), Inches(2.2 + row * 2.35), Inches(5.55), Inches(2.1), title, body, accent)
add_sn(sl, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — STORAGE SECTION
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_section_slide(sl, '05', '存储：大模型的「记忆」基石',
                  'HBM 堆叠竞赛 + NAND 闪存的大容量时代')
add_sn(sl, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — STORAGE DETAIL
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# 存储 · 记忆基石', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   'HBM4 即将登场，存内计算从概念走向量产', 36, NAVY, True)

mem_data = [
    ['', 'HBM3', 'HBM3e', 'HBM4', 'LPDDR5x (边缘)'],
    ['单堆容量', '24 GB', '36 GB', '48-64 GB', '16-32 GB'],
    ['带宽', '819 GB/s', '1.2 TB/s', '1.6-2.0 TB/s', '68 GB/s'],
    ['堆叠层数', '12-Hi', '12-Hi', '16-Hi', '—'],
    ['IO 位宽', '1024-bit', '1024-bit', '2048-bit', '64-bit'],
    ['主要供应商', 'SK海力士 / 三星', 'SK海力士 / 三星 / 美光',
     'SK海力士（先行）', '三星 / 美光'],
]
add_table(sl, Inches(1), Inches(2.2), Inches(11.3), mem_data,
          [Inches(1.8), Inches(2.375), Inches(2.375), Inches(2.375), Inches(2.375)])

add_card(sl, Inches(1), Inches(4.9), Inches(5.5), Inches(1.35),
         '\U0001f53a HBM 市场格局',
         'SK 海力士以 53% 份额领先（HBM3e 独家供货 NVIDIA），'
         '三星 38%，美光 9%。2026 全球 HBM 市场规模预计突破 400 亿美元。')
add_card(sl, Inches(6.8), Inches(4.9), Inches(5.5), Inches(1.35),
         '\U0001f53a 企业级 SSD',
         'PCIe 5.0 NVMe SSD 单盘 64TB。QLC NAND + 存算一体降低数据搬移能耗。'
         'AI 推理缓存池化需求拉动大容量 SSD。')

add_insight(sl, Inches(1), Inches(6.45), Inches(11.5), Inches(0.75),
            '\U0001f4a1 关键趋势',
            'HBM 是 GPU 性能释放的最大瓶颈之一。'
            'HBM4 通过 2048-bit 位宽和 16-Hi 堆叠将带宽翻倍。'
            'PIM（存内计算）让 HBM 不只是「存储」，还能执行矩阵运算。')
add_sn(sl, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — COOLING
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# 散热 · 能效挑战', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   '单机柜 100kW+ 时代，液冷从「可选项」变成「必选项」', 36, NAVY, True)

add_card(sl, Inches(1), Inches(2.2), Inches(5.5), Inches(2.2),
         '\U0001f321️ 功耗持续攀升',
         'GB300 NVL72 单机柜功耗高达 132kW。B200 单芯片 TDP 1000W，'
         '下一代 Rubin 预计 1500W+。传统风冷极限 ~30kW/机柜已被远远超越。',
         '132kW/rack')
add_card(sl, Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.2),
         '\U0001f4a7 液冷三大路线',
         '冷板式：成熟方案，PUE ~1.1，NVL72 标配。'
         '浸没式：单相/两相，PUE < 1.05，3M Novec 等氟化液成本仍高。'
         '芯片级：微通道直触，台积电 / Cooltera 研发中。',
         'PUE 1.05~1.15')

add_card(sl, Inches(1), Inches(4.6), Inches(3.55), Inches(1.6),
         '冷板液冷',
         '2026 年 AI 数据中心标配。CDU 单台 1.5MW 散热能力。'
         'Vertiv / CoolIT / Boyd 为主要供应商。')
add_card(sl, Inches(4.85), Inches(4.6), Inches(3.55), Inches(1.6),
         '浸没式液冷',
         'Tank + 氟化液方案。运维复杂度高，但散热密度极致。'
         '阿里 / 微软部分 AI 集群已部署。')
add_card(sl, Inches(8.7), Inches(4.6), Inches(3.55), Inches(1.6),
         '液冷产业链',
         '冷板 → CDU → 液冷管路 → 冷却塔。'
         '2026 全球液冷市场 ~120 亿美元（CAGR 35%+）。')
add_sn(sl, 13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — MARKET
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# 市场 · 投资格局', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   '2026 年 AI 基础设施市场规模', 36, NAVY, True)

add_stat_row(sl, Inches(2.1), [
    ('$420B', 'AI 芯片总市场', True),
    ('$180B', 'GPU / AI 加速器', False),
    ('$85B', 'HBM 存储', False),
    ('$28B', '光模块', False),
])

add_card(sl, Inches(1), Inches(3.6), Inches(3.55), Inches(2.0),
         '\U0001f3c6 芯片四强格局',
         'NVIDIA 以 ~80% AI GPU 份额绝对领先。AMD 追赶至 ~12%。'
         '自研 ASIC（Google TPU / AWS Trainium）合计 ~5%。'
         '华为 Ascend 在中国市场份额领先。')
add_card(sl, Inches(4.85), Inches(3.6), Inches(3.55), Inches(2.0),
         '\U0001f30f 国产替代加速',
         '华为 Ascend 910C 量产，寒武纪思元 690 进入数据中心。'
         '海光 DCU 兼容 ROCm 生态。'
         '中国 AI 芯片自给率 2026 年约 15-20%。')
add_card(sl, Inches(8.7), Inches(3.6), Inches(3.55), Inches(2.0),
         '\U0001f4c8 Capex 竞赛',
         '微软 / 谷歌 / 亚马逊 / Meta 2026 AI Capex 合计超 $300B。'
         'xAI Colossus 10 万 GPU 集群已投产，全球 AI 超算竞赛白热化。')
add_sn(sl, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — OUTLOOK
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
mono(sl, Inches(1), Inches(0.5), Inches(10), Inches(0.5), '# 展望 · 趋势总结', 13, BLUE, True)
tb(sl, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
   '2026-2028：六个确定性趋势', 36, NAVY, True)

trends = [
    ('\U0001f53a 1. 算力持续翻倍',
     'GPU 每 2 年算力翻倍（Blackwell → Rubin → Vera）。'
     'ASIC 推理芯片将占推理市场 30%+。'),
    ('\U0001f53a 2. 互联速率倍增',
     'Scale-Up：NVLink 6→7，UBB 2.0。Scale-Out：800G→1.6T→3.2T。'
     'CPO 将光引擎集成至交换机芯片，延迟降至 ns 级。'),
    ('\U0001f53a 3. 存储带宽破墙',
     'HBM4 → HBM4e，PIM 存内计算普及。'
     'CXL 3.0 内存池化打破「内存墙」。QLC NAND 256TB SSD。'),
    ('\U0001f53a 4. 液冷全面渗透',
     '2028 年 80%+ 新建 AI DC 采用液冷。PUE 从 1.4 → 1.05。'
     '液冷产业链 CAGR 35%+。'),
    ('\U0001f53a 5. 端侧 AI 崛起',
     'AI PC（NPU > 40 TOPS）、AI 手机、具身智能（机器人芯片）'
     '将创造新增长极。'),
    ('\U0001f53a 6. 地缘重塑供应链',
     'Chiplet + 先进封装（CoWoS / EMIB）成为「新摩尔定律」。'
     '东南亚封装测试产能扩建。'),
]
for i, (title, body) in enumerate(trends):
    col, row = i % 2, i // 2
    add_card(sl, Inches(1 + col * 5.85), Inches(2.2 + row * 1.55), Inches(5.55), Inches(1.3), title, body)
add_sn(sl, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Q&A
# ═══════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_bg(sl, WHITE)
# Grid
for i in range(0, 40):
    x = Emu(i * 40 * 12700)
    add_rect(sl, x, Inches(0), Pt(0.5), H, fill_color=GRID_C)
    y = Emu(i * 40 * 12700)
    add_rect(sl, Inches(0), y, W, Pt(0.5), fill_color=GRID_C)

center_tb(sl, Inches(0.5), Inches(3.5), '?', 150, GRID_C, True, 'Consolas', PP_ALIGN.CENTER)
center_tb(sl, Inches(3.4), Inches(1.2), '感谢聆听 · 欢迎交流', 46, NAVY, True)
center_tb(sl, Inches(4.3), Inches(0.9),
          'AI 基础设施正经历百年一遇的架构变革，'
          '每一个环节都值得深入跟踪。',
          16, DARK)

mono(sl, Inches(3.5), Inches(5.2), Inches(3), Inches(0.45),
     '\U0001f4e7  ai-research@example.com', 12, GRAY, align=PP_ALIGN.CENTER)
mono(sl, Inches(6), Inches(5.2), Inches(3), Inches(0.45),
     '\U0001f4ca  数据截止 2026.Q2', 12, GRAY, align=PP_ALIGN.CENTER)
center_tb(sl, Inches(5.7), Inches(0.5), 'Thank You · Questions & Discussion', 13, GRAY, font_name='Consolas')
add_sn(sl, 16)

# ── Save ──
output_path = r'd:\test\ai-industry-trends\AI产业发展趋势.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
